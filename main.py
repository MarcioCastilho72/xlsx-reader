from fastapi import FastAPI, File, UploadFile
from fastapi.responses import JSONResponse
from pptx import Presentation
from openpyxl import load_workbook
from io import BytesIO

app = FastAPI()

@app.post("/extract-text/")
async def extract_text(file: UploadFile = File(...)):
    if not file.filename.endswith(".pptx"):
        return JSONResponse(status_code=400, content={"error": "Only .pptx files are supported."})
    contents = await file.read()
    prs = Presentation(BytesIO(contents))
    full_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"
    return {"conteudoExtraido": full_text.strip()}


@app.post("/extract-xlsx/")
async def extract_xlsx(file: UploadFile = File(...)):
    if not file.filename.endswith(".xlsx"):
        return JSONResponse(status_code=400, content={"error": "Only .xlsx files are supported."})
    contents = await file.read()
    wb = load_workbook(filename=BytesIO(contents), data_only=True)
    full_text = ""
    for sheet in wb.worksheets:
        full_text += f"--- {sheet.title} ---\n"
        for row in sheet.iter_rows(values_only=True):
            full_text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
    return {"conteudoExtraido": full_text.strip()}